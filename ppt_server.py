from fastapi import FastAPI
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
import os

app = FastAPI()

class SlideData(BaseModel):
    title: str
    content: str

@app.post("/generate")
def generate_ppt(data: SlideData):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = data.title
    tf = body.text_frame
    for line in data.content.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
    os.makedirs("output", exist_ok=True)
    filepath = f"output/{data.title}.pptx"
    prs.save(filepath)
    return {"status": "ok", "path": filepath}
